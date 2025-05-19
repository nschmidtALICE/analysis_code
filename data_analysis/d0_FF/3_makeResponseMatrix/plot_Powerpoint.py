#
# How to run: python plot_Powerpoint.py -r Psi2S -p 5_10 -g True -z True
# How to run: python plot_Powerpoint.py -r Psi2S -g True -z True
#

from pptx import Presentation
from pptx.util import Inches, Pt
import argparse
import itertools
import os
import plot_Graphs
# ROOT
import ROOT

class PlotPowerpointObject:
    def __init__(self, resonance, ptRange, iszT,baseDir):



      self.binType = "UnBinned"
      self.binType2 = "Unbinned"
      print("Plot pptx for UnBinned fit")

      print("Plot pptx for {} GeV jets".format(ptRange))
      self.resonance = resonance
      self.ptRange   = ptRange
      if iszT=="True" or iszT=="TRUE":
        self.obsTag    = "zT"
      else:
        self.obsTag    = "dR"

      self.baseDir = baseDir
      self.plotDir = "{}/{}/".format(self.baseDir,ptRange)

    #- - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
    def makePowerpoint(self):
      

      qa_author = "Nicolas Schmidt, Los Alamos National lab"
      print("Author: %s" % qa_author)

      #--------------------------------------------------------------------
      # Start creating some introductory slides
      # Create a blank presentation
      qa = Presentation()
      
      # Define some slide layouts
      title_slide_layout  = qa.slide_layouts[0]
      blank_slide_layout  = qa.slide_layouts[6]
      
      # Make a title slide
      slide  = qa.slides.add_slide(title_slide_layout)
      title  = slide.shapes.title
      author = slide.placeholders[1]
      title.text  = "QA -- Fit of Mass and Time spectra"
      author.text = str(qa_author)

      #--------------------------------------------------------------------
      # Create the slides for all runs combined (Run = -1)
      # layout depending on the settings
      #plotRun(qa, blank_slide_layout, plotDir, isPtHard, isPbPb, isIncludePhos, -1, 0)
      self.plotTransitionSlide(qa, blank_slide_layout, "Mass fit evolution")
      self.plotMassFitVariationSlide(qa, blank_slide_layout)
      #-
      self.plotTransitionSlide(qa, blank_slide_layout, "Mass fit QA")
      self.plotMassFitSlidesQA(qa, blank_slide_layout)
      #-
      # self.plotTransitionSlide(qa, blank_slide_layout, "Time fit QA")
      # self.plotTimeFitSlidesQA(qa, blank_slide_layout)
      #-
      self.plotTransitionSlide(qa, blank_slide_layout, "Final Results")
      self.plotResultsSlide(qa, blank_slide_layout)
      #-
      

      #--------------------------------------------------------------------
      # Save the final presentation
      pptxName= self.baseDir + "/Fit_QA_{}_pt{}.pptx".format(self.resonance,self.ptRange)
      qa.save(pptxName)
      print("Saved presentation in: " +  pptxName)



    ###################################################################################################
    # Make slides of Mass fits parameters
    ###################################################################################################
    def plotTransitionSlide(self, qa, blank_slide_layout, titleName):

      titleleft   = Inches(2.)
      titletop    = Inches(3.2)
      titleheight = Inches(2.5)
      titlewidth  = Inches(7.)
     
      slide = qa.slides.add_slide(blank_slide_layout)
      txBox = slide.shapes.add_textbox(titleleft, titletop, titlewidth, titleheight)
      tf = txBox.text_frame
      tf.text = titleName
      set_text_properties(tf, 'Arial', 60, bold=True)

    ###################################################################################################
    # Make slides for Mass Fit QA
    ###################################################################################################
    def plotMassFitSlidesQA(self, qa, blank_slide_layout):
     
      fittype = "DCBFixed{}".format(self.binType2)
      #fittype = "DCBBinned"
      #fittype = "noSigBinned"
      #fittype = "noSigUnBinned"
      # self.plotSlideM(qa, blank_slide_layout, 0, fittype, "MParam_pol0R{}".format(self.binType))
      # self.plotSlideM(qa, blank_slide_layout, 1, fittype, "MParam_pol0R{}".format(self.binType))
      
      # self.plotSlideM(qa, blank_slide_layout, 0, fittype, "MParam_pol1R{}".format(self.binType))
      # self.plotSlideM(qa, blank_slide_layout, 1, fittype, "MParam_pol1R{}".format(self.binType))

      #Signal

      self.plotSlideM(qa, blank_slide_layout, 0, fittype, "MParam_MeanR{}".format(self.binType))
      self.plotSlideM(qa, blank_slide_layout, 1, fittype, "MParam_MeanR{}".format(self.binType))
      self.plotSlideM(qa, blank_slide_layout, 0, fittype, "MParam_Sig1R{}".format(self.binType))
      self.plotSlideM(qa, blank_slide_layout, 1, fittype, "MParam_Sig1R{}".format(self.binType))
      self.plotSlideM(qa, blank_slide_layout, 0, fittype, "MParam_Sig2R{}".format(self.binType))
      self.plotSlideM(qa, blank_slide_layout, 1, fittype, "MParam_Sig2R{}".format(self.binType))

      self.plotSlideM(qa, blank_slide_layout, 0, fittype, "MParam_CBFracR{}".format(self.binType))
      self.plotSlideM(qa, blank_slide_layout, 1, fittype, "MParam_CBFracR{}".format(self.binType))

      self.plotSlideM(qa, blank_slide_layout, 0, fittype, "MParam_alphaR{}".format(self.binType))
      self.plotSlideM(qa, blank_slide_layout, 1, fittype, "MParam_alphaR{}".format(self.binType))
      
      self.plotSlideM(qa, blank_slide_layout, 0, fittype, "MParam_NR{}".format(self.binType))
      self.plotSlideM(qa, blank_slide_layout, 1, fittype, "MParam_NR{}".format(self.binType))

    ###################################################################################################
    # Make slides for Time fit QA
    ###################################################################################################
    def plotTimeFitSlidesQA(self, qa, blank_slide_layout):
      
      self.plotSlideT(qa, blank_slide_layout, True, "TParam_MeanPR{}".format(self.binType))
      self.plotSlideT(qa, blank_slide_layout, False, "TParam_MeanPR{}".format(self.binType))
      self.plotSlideT(qa, blank_slide_layout, True, "TParam_MeanNPR{}".format(self.binType))
      self.plotSlideT(qa, blank_slide_layout, False, "TParam_MeanNPR{}".format(self.binType))
      
      self.plotSlideT(qa, blank_slide_layout, True, "TParam_Sig1R{}".format(self.binType))
      self.plotSlideT(qa, blank_slide_layout, False, "TParam_Sig1R{}".format(self.binType))
      self.plotSlideT(qa, blank_slide_layout, True, "TParam_Sig2R{}".format(self.binType))
      self.plotSlideT(qa, blank_slide_layout, False, "TParam_Sig2R{}".format(self.binType))
      self.plotSlideT(qa, blank_slide_layout, True, "TParam_GaussFracR{}".format(self.binType))
      self.plotSlideT(qa, blank_slide_layout, False, "TParam_GaussFracR{}".format(self.binType))
      
      self.plotSlideT(qa, blank_slide_layout, True, "TParam_ExpDecayR{}".format(self.binType))
      self.plotSlideT(qa, blank_slide_layout, False, "TParam_ExpDecayR{}".format(self.binType))

      self.plotSlideT(qa, blank_slide_layout, True, "TParam_BDecayFracR{}".format(self.binType))
      self.plotSlideT(qa, blank_slide_layout, False, "TParam_BDecayFracR{}".format(self.binType))
      
      self.plotSlideT(qa, blank_slide_layout, True, "TParam_YieldR{}".format(self.binType))
      self.plotSlideT(qa, blank_slide_layout, False, "TParam_YieldR{}".format(self.binType))

    ###################################################################################################
    # Make slides for Mass Fit Variation
    ###################################################################################################
    def plotMassFitVariationSlide(self, qa, blank_slide_layout):

      maxBin=5

      for bin in range(0,maxBin):
        self.plotSlideMVariation(qa, blank_slide_layout, bin)

    ###################################################################################################
    # Make slides of Final results
    ###################################################################################################
    def plotResultsSlide(self, qa, blank_slide_layout):

      titleleft   = Inches(1.)
      titletop    = Inches(0.1)
      titleheight = Inches(0.8)
      titlewidth  = Inches(8.)
      
      slide = qa.slides.add_slide(blank_slide_layout)
      txBox = slide.shapes.add_textbox(titleleft, titletop, titlewidth, titleheight)
      # Title
      tf = txBox.text_frame
      tf.text = "Final Results for {} with pT(Jet)={} GeV".format(self.resonance,self.ptRange)
      set_text_properties(tf, 'Arial', 40, bold=True)


      imgWidth = 5.0
      shift = 1.6
      # columns
      leftEdge1  = 0.1
      leftEdge2  = imgWidth
      # rows
      topEdge1  = 1.1+shift
      
      topB    = Inches(2.3)
      heightB = Inches(1)
      widthB  = Inches(4)
      left1B  = Inches(1.)
      left2B  = Inches(6.)
      txBox1 = slide.shapes.add_textbox(left1B, topB, widthB, heightB)
      tf1 = txBox1.text_frame
      tf1.text = "Absolute Yield"
      set_text_properties(tf1, 'Arial', 20)

      img_path = self.plotDir + "YieldAbs{}_All.png".format(self.binType)
      left  = Inches(leftEdge1)
      top   = Inches(topEdge1)
      width = Inches(imgWidth)
      slide.shapes.add_picture(img_path, left, top, width=width)
     
      txBox2 = slide.shapes.add_textbox(left2B, topB, widthB, heightB)
      tf2 = txBox2.text_frame
      tf2.text = "dN/dz_T"
      set_text_properties(tf2, 'Arial', 20)

      img_path = self.plotDir + "YieldPer_{}{}_All.png".format(self.obsTag,self.binType)#Eli
      left  = Inches(leftEdge2)
      top   = Inches(topEdge1)
      width = Inches(imgWidth)
      slide.shapes.add_picture(img_path, left, top, width=width)

    ###################################################################################################
    # Plot slides with track information
    ###################################################################################################
    def plotSlideMVariation(self, qa, blank_slide_layout, Bin):

      titleleft   = Inches(2.6)
      titletop    = Inches(0.)
      titleheight = Inches(0.8)
      titlewidth  = Inches(5.)
      slide = qa.slides.add_slide(blank_slide_layout)
      txBox = slide.shapes.add_textbox(titleleft, titletop+0.2, titlewidth, titleheight)


      imgWidth = 5.1
      shift = 1.6
      # columns
      leftEdge1  = 0.1
      leftEdge2  = imgWidth
      #leftEdge3  = 6.6
      # rows
      topEdge1  = 1.1+shift
      topEdge2  = 3.5+shift
      #lowEdge     = 4.6
      
      # Title
      tf = txBox.text_frame
      tf.text = "Mass fits"
      set_text_properties(tf, 'Arial', 40, bold=True)

      #Comparisions of yields
      img_path = self.plotDir + "YieldAbs{}_All.png".format(self.binType)
      print("Add picture" + img_path)
      left  = Inches(5)
      top   = Inches(0.05)
      width = Inches(3.8)
      slide.shapes.add_picture(img_path, left, top, width=width)

      #-
      img_path = self.plotDir + "MassFits_{}/Bin{}_SGauss{}.png".format(self.obsTag,Bin,self.binType2)
      left  = Inches(leftEdge1)
      top   = Inches(topEdge1)
      width = Inches(imgWidth)
      slide.shapes.add_picture(img_path, left, top, width=width)

      img_path = self.plotDir + "MassFits_{}/Bin{}_DGauss{}.png".format(self.obsTag,Bin,self.binType2)
      left  = Inches(leftEdge2)
      top   = Inches(topEdge1)
      width = Inches(imgWidth)
      slide.shapes.add_picture(img_path, left, top, width=width)

      img_path = self.plotDir + "MassFits_{}/Bin{}_DCB{}.png".format(self.obsTag,Bin,self.binType2)
      left  = Inches(leftEdge1)
      top   = Inches(topEdge2)
      width = Inches(imgWidth)
      slide.shapes.add_picture(img_path, left, top, width=width)
  
      img_path = self.plotDir + "MassFits_{}/Bin{}_DCBFixed{}.png".format(self.obsTag,Bin,self.binType2)
      left  = Inches(leftEdge2)
      top   = Inches(topEdge2)
      width = Inches(imgWidth)
      slide.shapes.add_picture(img_path, left, top, width=width)

    ###################################################################################################
    # Plot slides with track information
    ###################################################################################################
    def plotSlideM(self, qa, blank_slide_layout, binPortion, fittype, paramName):

      titleleft   = Inches(2.6)
      titletop    = Inches(0.)
      titleheight = Inches(0.8)
      titlewidth  = Inches(5.)
      slide = qa.slides.add_slide(blank_slide_layout)
      txBox = slide.shapes.add_textbox(titleleft, titletop+0.2, titlewidth, titleheight)
      # Title
      tf = txBox.text_frame
      tf.text = "Mass fits"
      set_text_properties(tf, 'Arial', 40, bold=True)

      #-Fit identifier
      topB    = Inches(2.3)
      heightB = Inches(1)
      widthB  = Inches(4)
      left1B  = Inches(1.)
      left2B  = Inches(6.)
      txBox1 = slide.shapes.add_textbox(left1B, topB, widthB, heightB)
      tf1 = txBox1.text_frame
      #tf1.text = "OldFit T_FuncFit"
      #tf1.text = "NewFit Cuts"
      tf1.text = "March Data - most Cuts"
      #tf1.text = "March Data - Cuts"
      #tf1.text = "March Data - no Cuts"
      #tf1.text = "Sept Data - new Fit"
      set_text_properties(tf1, 'Arial', 20)


      imgWidth = 5.1
      shift = 1.6
      # columns
      leftEdge1  = 0.1
      leftEdge2  = imgWidth
      #leftEdge3  = 6.6
      # rows
      topEdge1  = 1.1+shift
      topEdge2  = 3.5+shift
      #lowEdge     = 4.6
      
      img_path = self.plotDir + paramName+ ".png"
      print("Add picture" + img_path)
      left  = Inches(5)
      top   = Inches(0.05)
      width = Inches(3.8)
      slide.shapes.add_picture(img_path, left, top, width=width)

      #-
      if binPortion==0:
        #if it is a Pt hard bin production plot more info on the slides
        img_path = self.plotDir + "MassFits_{}/Bin0_{}.png".format(self.obsTag,fittype)
        left  = Inches(leftEdge1)
        top   = Inches(topEdge1)
        width = Inches(imgWidth)
        slide.shapes.add_picture(img_path, left, top, width=width)

        img_path = self.plotDir + "MassFits_{}/Bin1_{}.png".format(self.obsTag,fittype)
        left  = Inches(leftEdge2)
        top   = Inches(topEdge1)
        width = Inches(imgWidth)
        slide.shapes.add_picture(img_path, left, top, width=width)

        img_path = self.plotDir + "MassFits_{}/Bin2_{}.png".format(self.obsTag,fittype)
        left  = Inches(leftEdge1)
        top   = Inches(topEdge2)
        width = Inches(imgWidth)
        slide.shapes.add_picture(img_path, left, top, width=width)

        img_path = self.plotDir + "MassFits_{}/Bin3_{}.png".format(self.obsTag,fittype)
        left  = Inches(leftEdge2)
        top   = Inches(topEdge2)
        width = Inches(imgWidth)
        slide.shapes.add_picture(img_path, left, top, width=width)

      elif binPortion==1:
        #if it is a Pt hard bin production plot more info on the slides
        img_path = self.plotDir + "MassFits_{}/Bin4_{}.png".format(self.obsTag,fittype)
        left  = Inches(leftEdge1)
        top   = Inches(topEdge1)
        width = Inches(imgWidth)
        slide.shapes.add_picture(img_path, left, top, width=width)

        img_path = self.plotDir + "MassFits_{}/Bin5_{}.png".format(self.obsTag,fittype)
        left  = Inches(leftEdge2)
        top   = Inches(topEdge1)
        width = Inches(imgWidth)
        slide.shapes.add_picture(img_path, left, top, width=width)

        # img_path = self.plotDir + "MassFits_{}/Bin6_{}.png".format(self.obsTag,fittype)
        # left  = Inches(leftEdge1)
        # top   = Inches(topEdge2)
        # width = Inches(imgWidth)
        # slide.shapes.add_picture(img_path, left, top, width=width)


      elif binPortion==2:
            #if it is a Pt hard bin production plot more info on the slides
            img_path = self.plotDir + "MassFits_{}/Bin8_{}.png".format(self.obsTag,fittype)
            left  = Inches(leftEdge1)
            top   = Inches(topEdge1)
            width = Inches(imgWidth)
            slide.shapes.add_picture(img_path, left, top, width=width)

            img_path = self.plotDir + "MassFits_{}/Bin9_{}.png".format(self.obsTag,fittype)
            left  = Inches(leftEdge2)
            top   = Inches(topEdge1)
            width = Inches(imgWidth)
            slide.shapes.add_picture(img_path, left, top, width=width)

    ###################################################################################################
    # Plot slides with track information
    ###################################################################################################
    def plotSlideT(self, qa, blank_slide_layout, FirstbinPortion, paramName):

      titleleft   = Inches(2.6)
      titletop    = Inches(0.)
      titleheight = Inches(0.8)
      titlewidth  = Inches(5.)
      slide = qa.slides.add_slide(blank_slide_layout)
      txBox = slide.shapes.add_textbox(titleleft, titletop+0.2, titlewidth, titleheight)


      imgWidth = 5.1
      shift = 1.6
      # columns
      leftEdge1  = 0.1
      leftEdge2  = imgWidth
      #leftEdge3  = 6.6
      # rows
      topEdge1  = 1.1+shift
      topEdge2  = 3.5+shift
      #lowEdge     = 4.6
      
      
      # Title
      tf = txBox.text_frame
      tf.text = "Time fits"
      set_text_properties(tf, 'Arial', 40, bold=True)

      img_path = self.plotDir + paramName+ ".png"
      print("Add picture" + img_path)
      left  = Inches(5)
      top   = Inches(0.05)
      width = Inches(3.8)
      slide.shapes.add_picture(img_path, left, top, width=width)

      #-
      if FirstbinPortion:
        #if it is a Pt hard bin production plot more info on the slides
        img_path = self.plotDir + "TimeFits_{}/LifetimeFitsFun_0.png".format(self.obsTag)
        left  = Inches(leftEdge1)
        top   = Inches(topEdge1)
        width = Inches(imgWidth)
        slide.shapes.add_picture(img_path, left, top, width=width)

        img_path = self.plotDir + "TimeFits_{}/LifetimeFitsFun_1.png".format(self.obsTag)
        left  = Inches(leftEdge2)
        top   = Inches(topEdge1)
        width = Inches(imgWidth)
        slide.shapes.add_picture(img_path, left, top, width=width)

        img_path = self.plotDir + "TimeFits_{}/LifetimeFitsFun_2.png".format(self.obsTag)
        left  = Inches(leftEdge1)
        top   = Inches(topEdge2)
        width = Inches(imgWidth)
        slide.shapes.add_picture(img_path, left, top, width=width)

        img_path = self.plotDir + "TimeFits_{}/LifetimeFitsFun_3.png".format(self.obsTag)
        left  = Inches(leftEdge2)
        top   = Inches(topEdge2)
        width = Inches(imgWidth)
        slide.shapes.add_picture(img_path, left, top, width=width)

      else:
        #if it is a Pt hard bin production plot more info on the slides
        img_path = self.plotDir + "TimeFits_{}/LifetimeFitsFun_4.png".format(self.obsTag)
        left  = Inches(leftEdge1)
        top   = Inches(topEdge1)
        width = Inches(imgWidth)
        slide.shapes.add_picture(img_path, left, top, width=width)

        img_path = self.plotDir + "TimeFits_{}/LifetimeFitsFun_5.png".format(self.obsTag)
        left  = Inches(leftEdge2)
        top   = Inches(topEdge1)
        width = Inches(imgWidth)
        slide.shapes.add_picture(img_path, left, top, width=width)
        
        img_path = self.plotDir + "TimeFits_{}/LifetimeFitsFun_6.png".format(self.obsTag)
        left  = Inches(leftEdge1)
        top   = Inches(topEdge2)
        width = Inches(imgWidth)
        slide.shapes.add_picture(img_path, left, top, width=width)

def set_text_properties(text_frame, font_name, font_size, bold=False, italic=False):
    """
    Set text properties directly without using fit_text (which doesn't work on all Linux systems)
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    
    # Set properties for all paragraphs in the text frame
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = 1  # PP_ALIGN.CENTER
        for run in paragraph.runs:
            font = run.font
            font.name = font_name
            font.size = Pt(font_size)
            font.bold = bold
            font.italic = italic
            # Optional: set a color
            font.color.rgb = RGBColor(0, 0, 0)  # Black text

#- - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
def plotPowerpoint(resKey, ptRange, plotGraphs, iszT):
 
  if "D0" in resKey:
    print("Plot pptx for D0")
    pTRangeArray = ["5_10","10_15","15_20","20_30"]
    #pTRangeArray = ["5_50"]
    # print("Error not the correct resonance key - abort")
    # return

  baseDir = "/media/niviths/local/analysis_code/data_analysis/d0_FF/2_fitData/D0_FF"

  if ptRange == "":
     for range in pTRangeArray:
          if plotGraphs=="True":
            print("plot the fit parameter QA plots first")
            plot_Graphs.PlotGraphs(resKey, range, iszT,baseDir)
          
          graphs = PlotPowerpointObject(resKey, range, iszT, baseDir)
          graphs.makePowerpoint()
  else:
      if plotGraphs=="True":
        print("plot the fit parameter QA plots first")
        plot_Graphs.PlotGraphs(resKey, ptRange, iszT,baseDir)
      
      graphs = PlotPowerpointObject(resKey, ptRange,iszT,baseDir)
      graphs.makePowerpoint()

#---------------------------------------------------------------------------------------------------
if __name__ == '__main__':
  # Define arguments def plotPowerpoint(resKey, ptRange):
  parser = argparse.ArgumentParser(description="Create an overview of the fit parameters for the mass and time fits")
  parser.add_argument("-r", "--resKey", action="store",
                      type=str, metavar="resKey",
                      default="D0",
                      help="List of runs to iterate over")
  parser.add_argument("-p", "--ptRange", action="store",
                      type=str, metavar="ptRange",
                      default="",
                      help="pt range for the jet")
  parser.add_argument("-g", "--plotGraphs", action="store",
                      type=str, metavar="plotGraphs",
                      #metavar="plotGraphs",
                      #default=False,
                      help="should the graphs be created too?")
  parser.add_argument("-z", "--iszT", action="store",
                      type=str, metavar="iszT",
                      default="True",
                      #metavar="isBinned",
                      help="is it for the zT observable")
  # Parse the arguments
  args = parser.parse_args()
  
  plotPowerpoint(args.resKey, args.ptRange, args.plotGraphs, args.iszT)
# How to run: python plot_Powerpoint.py -r D0 -p 5_10 -g True -z True
# How to run: python plot_Powerpoint.py -r D0 -g True -z True
# How to run: python plot_Powerpoint.py -r D0 -g True -z True
